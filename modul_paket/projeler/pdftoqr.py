from pptx import Presentation
import qrcode

def pptx_to_qrcodes(pptx_path, output_directory):
    presentation = Presentation(pptx_path)

    for slide_number, slide in enumerate(presentation.slides):
        text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"

        # QR kodu oluştur
        qr = qrcode.QRCode(
            version=1,
            error_correction=qrcode.constants.ERROR_CORRECT_L,
            box_size=10,
            border=4,
        )
        qr.add_data(text)
        qr.make(fit=True)

        # QR kodunu resim olarak kaydet
        img = qr.make_image(fill_color="black", back_color="white")
        output_path = f"{output_directory}/slide_{slide_number + 1}_qr_code.png"
        img.save(output_path)

if __name__ == "__main__":
    pptx_file_path = "C:\\Users\\tahak\\Desktop\\Yeni klasör (4)\\yarabakim.pptx"
    output_directory = "C:\\Users\\tahak\\Desktop\\Yeni klasör (4)"

    pptx_to_qrcodes(pptx_file_path, output_directory)
